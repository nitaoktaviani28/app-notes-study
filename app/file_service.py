from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation


def ensure_folder(path: str) -> None:
    Path(path).mkdir(parents=True, exist_ok=True)


def extract_text_from_material(file_path: str) -> str:
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext == ".pdf":
        reader = PdfReader(str(path))
        return "\n".join((page.extract_text() or "") for page in reader.pages)

    if ext == ".pptx":
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts)

    return ""
